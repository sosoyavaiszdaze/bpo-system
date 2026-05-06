"""scripts/generate_proposal_pptx.py の単体テスト（Day 5.3 Task 4）

カバー範囲:
    (a) 全 13 ページが生成されること
    (b) ページ8の Top5 統合表が priority_ranker / impact_estimator の出力と数値一致すること
    (c) スピーカーノートに ADR 番号が機械的に挿入されていること
    (d) ファイル名規則（日付・バージョン）が守られていること
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

# 5/7 P1-A: python-pptx は提案資料生成専用のオプショナル依存。
# 未インストール環境 (CI 一部 / requirements.lock 不在環境) では
# import 時点でモジュールごと collect-skip させ、CI の "ModuleNotFoundError"
# による全件 fail を防ぐ。pip install python-pptx で全テスト通る。
pytest.importorskip("pptx", reason="python-pptx 未インストール (オプショナル依存)")

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from scripts.generate_proposal_pptx import (
    determine_output_path,
    gather_top5_with_3layer,
    generate,
    get_client_config,
    load_audit_results,
    load_clients_yaml,
)


@pytest.fixture(scope="module")
def generated_pptx(tmp_path_factory):
    """1 度だけ pptx を生成し、複数テストで再利用"""
    out_dir = tmp_path_factory.mktemp("pptx_out")
    out_path = generate("pilotton", output_dir=out_dir)
    return out_path


def _open_pptx(path: Path):
    """python-pptx で pptx を開く"""
    from pptx import Presentation
    return Presentation(str(path))


# =============================================================================
# (a) 全 13 ページが生成されること
# =============================================================================
def test_pptx_has_13_pages(generated_pptx):
    prs = _open_pptx(generated_pptx)
    assert len(prs.slides) == 13, f"期待 13 ページ、実際 {len(prs.slides)} ページ"


def test_pptx_is_widescreen_16_9(generated_pptx):
    prs = _open_pptx(generated_pptx)
    # 16:9 サイズ確認（cm 表現は浮動小数のため tolerance を持たせる）
    assert abs(prs.slide_width.cm - 33.867) < 0.1
    assert abs(prs.slide_height.cm - 19.05) < 0.1


# =============================================================================
# (b) ページ8の Top5 統合表が runtime データと数値一致
# =============================================================================
def test_p8_top5_numerical_alignment(generated_pptx):
    """ページ 8（Top5 統合表）の数値が impact_estimator の出力と完全一致することを確認"""
    audit_data = load_audit_results("pilotton")
    audit = audit_data.get("ads_audit", {})
    clients_yaml = load_clients_yaml()
    client_cfg = get_client_config("pilotton", clients_yaml)
    runtime = gather_top5_with_3layer(audit, client_cfg)

    # 期待: priority_ranker 実出力の各 rule_id, conservative_yen が pptx P8 に含まれる
    prs = _open_pptx(generated_pptx)
    p8 = prs.slides[7]  # 0-indexed
    p8_text = ""
    for shape in p8.shapes:
        if shape.has_text_frame:
            p8_text += shape.text_frame.text + "\n"

    for row in runtime["top5_rows"]:
        # rule_id が含まれている
        assert row["rule_id"] in p8_text, f"P8 に {row['rule_id']} が含まれていない"
        # 確実値が円表記で含まれている
        cons_str = f"¥{row['conservative']:,}"
        assert cons_str in p8_text, f"P8 に {row['rule_id']} の確実値 {cons_str} が含まれていない"

    # 合計行も確認
    total_cons = f"¥{runtime['total']['conservative']:,}"
    assert total_cons in p8_text, f"P8 に Top5 合計確実値 {total_cons} が含まれていない"


def test_top5_order_matches_priority_ranker(generated_pptx):
    """Top5 順序が priority_ranker の動的出力に追従していることを確認"""
    audit_data = load_audit_results("pilotton")
    clients_yaml = load_clients_yaml()
    client_cfg = get_client_config("pilotton", clients_yaml)
    runtime = gather_top5_with_3layer(audit_data["ads_audit"], client_cfg)

    rule_ids_runtime = [r["rule_id"] for r in runtime["top5_rows"]]
    assert len(rule_ids_runtime) == 5, f"Top5 が 5 件でない: {rule_ids_runtime}"

    # P8 内の出現位置で順序を検証
    prs = _open_pptx(generated_pptx)
    p8 = prs.slides[7]
    full_text = ""
    for shape in p8.shapes:
        if shape.has_text_frame:
            full_text += shape.text_frame.text + "\n"

    positions = []
    for rid in rule_ids_runtime:
        pos = full_text.find(f"[{rid}]")
        assert pos >= 0, f"P8 に {rid} のラベルが見つからない"
        positions.append(pos)

    # 順序が単調増加していること
    assert positions == sorted(positions), (
        f"Top5 順序が priority_ranker 出力と一致しない。runtime={rule_ids_runtime}, "
        f"positions={positions}"
    )


# =============================================================================
# (c) スピーカーノートに ADR 番号が機械的に挿入されている
# =============================================================================
def test_speaker_notes_contain_adr_references(generated_pptx):
    """全スライドのスピーカーノートに【根拠ADR】行が含まれることを確認"""
    prs = _open_pptx(generated_pptx)
    slides_with_adr = 0
    slides_without_adr = []
    for i, slide in enumerate(prs.slides, 1):
        notes = slide.notes_slide.notes_text_frame.text
        if "【根拠ADR】" in notes:
            slides_with_adr += 1
        else:
            slides_without_adr.append(i)
    assert slides_with_adr == 13, (
        f"全 13 ページに ADR 引用が必要、欠落ページ: {slides_without_adr}"
    )


def test_speaker_notes_contain_data_source(generated_pptx):
    """全スライドのスピーカーノートに【データソース】行が含まれることを確認"""
    prs = _open_pptx(generated_pptx)
    for i, slide in enumerate(prs.slides, 1):
        notes = slide.notes_slide.notes_text_frame.text
        assert "【データソース】" in notes, f"P{i} のノートに【データソース】が無い"


def test_p8_notes_reference_adr_001_002_003(generated_pptx):
    """P8（最重要ページ）のスピーカーノートに ADR-001/002/003 の 3 本全てが引用されている"""
    prs = _open_pptx(generated_pptx)
    p8_notes = prs.slides[7].notes_slide.notes_text_frame.text
    for adr in ["ADR-001", "ADR-002", "ADR-003"]:
        assert adr in p8_notes, f"P8 ノートに {adr} の引用が無い"


# =============================================================================
# (d) ファイル名規則（日付・バージョン）
# =============================================================================
def test_filename_matches_versioning_pattern(generated_pptx):
    """生成ファイル名が <client>_proposal_v<N>.pptx 形式であることを確認"""
    import re
    pattern = re.compile(r"^pilotton_proposal_v\d+(?:_\d{4})?\.pptx$")
    assert pattern.match(generated_pptx.name), (
        f"ファイル名が規則違反: {generated_pptx.name}"
    )


def test_versioning_increments_for_existing_files(tmp_path):
    """既存ファイルがあれば自動でバージョン番号がインクリメントされることを確認"""
    # v1 を事前に作成
    v1_path = tmp_path / "pilotton_proposal_v1.pptx"
    v1_path.write_bytes(b"dummy")
    # determine_output_path が v2 を返すこと
    next_path = determine_output_path("pilotton", tmp_path)
    assert next_path.name == "pilotton_proposal_v2.pptx", (
        f"期待 v2、実際 {next_path.name}"
    )

    # v3 まで作成して次は v4 か確認
    (tmp_path / "pilotton_proposal_v2.pptx").write_bytes(b"dummy")
    (tmp_path / "pilotton_proposal_v3.pptx").write_bytes(b"dummy")
    next_path = determine_output_path("pilotton", tmp_path)
    assert next_path.name == "pilotton_proposal_v4.pptx"


# =============================================================================
# 補足: 数値ハードコード禁止の検証（grep ベース、簡易）
# =============================================================================
def test_no_hardcoded_pilotton_numbers_in_script():
    """generate_proposal_pptx.py 本体に Day 5.2 真値（CV 161 / CPA 8984 / 月次 1446455 等）が
    リテラルとして埋め込まれていないことを確認（SSoT 強制）。

    カンマ表記（¥8,984）の形でも検出できるよう、複数表現を確認する。
    凡例の業界平均値（4,500 など）は benchmarks.yaml の値で SSoT のため許容。
    """
    src = (PROJECT_ROOT / "scripts" / "generate_proposal_pptx.py").read_text(encoding="utf-8")

    # Day 5.2 真値の各種表記を網羅
    forbidden_patterns = [
        # CV 161（単独整数だと 16161 等にも誤マッチするので前後 quote/space で確認）
        ("CV161",       "CV 161 リテラル",   r"\bCV\s*161\b|161\s*件"),
        ("CPA8984",     "CPA ¥8,984 リテラル", r"¥\s*8[,，]?984|\b8984\b"),
        ("Spend14M",    "月次支出 ¥1,446,455 リテラル", r"¥\s*1[,，]?446[,，]?455|\b1446455\b"),
        ("CTR237",      "CTR 2.37% リテラル", r"\b2\.37\s*%"),
    ]
    import re
    truly_forbidden = []
    for tag, desc, pattern in forbidden_patterns:
        if re.search(pattern, src):
            truly_forbidden.append(f"{tag}: {desc}")
    assert not truly_forbidden, (
        f"スクリプト内に Day 5.2 真値のハードコードが検出されました:\n  - "
        + "\n  - ".join(truly_forbidden)
        + "\n全ての pilotton 固有数値は runtime（audit / impact_estimator / benchmarks）から取得すること"
    )
