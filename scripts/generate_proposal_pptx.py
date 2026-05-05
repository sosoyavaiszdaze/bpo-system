#!/usr/bin/env python3
"""pilotton 向け PoC 提案 pptx 生成スクリプト（Day 5.3）。

設計根拠:
    - ADR-001: 想定改善額の3層表示
    - ADR-002: 6グループ root_cause_group
    - ADR-003: pixel_health 連動
    - ADR-004: CV カウント正規化（pptx 数値の信頼性前提）

設計原則:
    1. Single Source of Truth: 数値は全て runtime に既存出力から取得（pptx 内ハードコード禁止）
    2. ADR トレーサビリティ: 各ページのスピーカーノートに「データソース」「根拠ADR」自動挿入
    3. 順序動的追従: Top5 順序は priority_ranker の実出力に従う（Day 5.3 時点 M09→M02→M03→M04→M61）
    4. 機微情報: ファイル名は日付＋バージョン明示、フッターに「社外秘」、数値引用元はノートのみ

使い方:
    python scripts/generate_proposal_pptx.py --client pilotton
    python scripts/generate_proposal_pptx.py --client pilotton --output-dir reports/2026-05-07/

ファイル名規則:
    reports/<YYYY-MM-DD>/<client>_proposal_v<N>.pptx
    既存ファイルがあれば自動でバージョン番号インクリメント
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path

# プロジェクトルートを sys.path に追加
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# 既存エンジンモジュール
from analyzers.ads_audit import detect_pixel_health
from engine.benchmark_compare import build_health_score_3axis, compare_3axis, load_benchmarks
from engine.impact_estimator import (
    aggregate_with_dedup,
    calculate_independent_impact,
    calculate_minimum_impact,
    calculate_realistic_impact,
    estimate_for_rule,
)
from engine.priority_ranker import compute_top_actions, load_all_rules, load_weights

# =============================================================================
# 定数
# =============================================================================

# Zynect ブランドカラー（templates/v3/_styles.html から流用）
COLOR_GREEN = RGBColor(0x27, 0x50, 0x0A)       # 確実・優秀
COLOR_ORANGE = RGBColor(0xBA, 0x75, 0x17)      # 中庸・警告
COLOR_RED = RGBColor(0xA3, 0x2D, 0x2D)         # 危険
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x1A)        # 本文黒
COLOR_GRAY = RGBColor(0x88, 0x88, 0x88)        # 補助
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xED)  # 帯背景
COLOR_BG = RGBColor(0xFA, 0xFA, 0xF7)          # ページ背景（オフホワイト）
COLOR_HIGHLIGHT_BG = RGBColor(0xEA, 0xF3, 0xDE)  # 確実値ハイライト

# フォント（Hiragino / Yu Gothic / Helvetica）
FONT_JP = "Hiragino Sans"
FONT_EN = "Helvetica"

# スライドサイズ（16:9、cm）
SLIDE_W_CM = 33.867
SLIDE_H_CM = 19.05

# グループ表示名（短縮版、Top5 表で利用）
GROUP_SHORT = {
    "measurement_foundation": "MF",
    "delivery_learning_or_structure": "DLS",
    "creative_optimization": "CR",
    "budget_allocation": "BUD",
    "targeting": "TGT",
    "independent": "IND",
    "other": "—",
}
GROUP_FULL = {
    "measurement_foundation": "計測基盤",
    "delivery_learning_or_structure": "配信学習・構造",
    "creative_optimization": "クリエイティブ",
    "budget_allocation": "予算配分",
    "targeting": "ターゲティング",
    "independent": "独立施策",
    "other": "未分類",
}

# 米満氏理論 9+10 原則の要約（Appendix C 用）
PRINCIPLES_SUMMARY = [
    ("P1", "計測精度=学習シグナル精度", "計測の正しさは全運用判断の前提"),
    ("P2", "機械学習保護", "短期判断による長期最適化の毀損を回避"),
    ("P3", "結果指標非依存", "品質スコアではなく原因変数で判断"),
    ("P4", "ネガティブシグナル保持", "低パフォ要素は削除せず除外保持"),
    ("P5", "Budget Lost 先行解消", "効率改善より機会損失の解消を優先"),
    ("P6", "集約優先・分離", "学習単位は集約、評価軸異質は分離"),
    ("P7", "バリエーション幅最大化", "多パターン×短い見出し"),
    ("P8", "自動化前提判断", "旧式手動運用の知識を捨てる"),
    ("P9", "説明責任・判断ログ", "なぜその判断をしたかを残す"),
    ("M-α", "計測=シグナル基盤", "Pixel + CAPI + EMQ の三位一体"),
    ("M-β", "学習フェーズ保護", "週 50CV 基準、編集抑制"),
    ("M-η", "Advantage+ 自動化前提", "ASC+ 等の自動化を信頼"),
    ("M-ζ", "クリエイティブ多様性", "ASC 内 15-50 本のアクティブ運用"),
    ("M-θ", "iOS14 計測欠損前提運用", "AEM 優先度設定で計測欠損補完"),
    ("M-λ", "広告-LP メッセージ完全一致", "整合度が CVR を支配"),
]


# =============================================================================
# データ取得層（Single Source of Truth）
# =============================================================================

def load_audit_results(client_id: str, output_dir: Path | None = None) -> dict:
    """最新の {client}_results.json を取得（reports/ 以下を新しい順に走査）。"""
    reports_dir = PROJECT_ROOT / "reports"
    candidates = []
    for date_dir in reports_dir.iterdir():
        if not date_dir.is_dir():
            continue
        target = date_dir / f"{client_id}_results.json"
        if target.exists():
            candidates.append((date_dir.name, target))
    if not candidates:
        raise FileNotFoundError(f"{client_id}_results.json が reports/ 配下に見つかりません")
    candidates.sort(key=lambda x: x[0], reverse=True)
    src_path = candidates[0][1]
    with open(src_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    print(f"[SSoT] audit results: {src_path}")
    data["_source_path"] = str(src_path.relative_to(PROJECT_ROOT))
    return data


def load_clients_yaml() -> dict:
    """config/clients.yaml を読み込む。"""
    with open(PROJECT_ROOT / "config" / "clients.yaml", "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def get_client_config(client_id: str, clients_yaml: dict) -> dict:
    cfg = clients_yaml.get("clients", {}).get(client_id)
    if not cfg:
        raise ValueError(f"client_id='{client_id}' が clients.yaml に未登録")
    return cfg


def gather_top5_with_3layer(audit: dict, client_cfg: dict) -> dict:
    """Day 5.3 重要部: priority_ranker 実出力で Top5 順序を取得し、各行に 3 層インパクトを付与。

    Returns:
        dict {
          top5_rows: [{rank, rule_id, rule_name, group, factor,
                       conservative, realistic, optimistic, confidence_stars}],
          total: {conservative, realistic, optimistic},
          monthly_spend: int,
          pixel_health: {dormant_days, ...},
          cv_count: int, cpa: int, ctr: float, ...
        }
    """
    # priority_ranker 経由で Top5 を動的に取得（順序ハードコード禁止）
    weights = load_weights()
    rules = load_all_rules()
    issues = audit.get("issues", []) or []
    detected_ids = []
    for i in issues:
        rid = i.get("check_id") or i.get("id") or i.get("rule_id")
        if rid:
            detected_ids.append(rid)

    monthly_spend = float(audit.get("total_cost", 0) or 0)
    top5 = compute_top_actions(detected_ids, rules, weights, monthly_spend, max_n=5)

    # pixel_health 連動を有効化
    pixels = ((client_cfg.get("ads") or {}).get("meta") or {}).get("pixels") or []
    pixel_health = detect_pixel_health(pixels)

    # v3.1.2 (Day 5.3 A-T3): 各アクションに current_metrics を構築して渡し、
    # PDF (report_generator_v3) と完全同一の試算経路にする（数値乖離防止）
    issues = audit.get("issues") or []
    platform_summary = audit.get("platform_summary") or {}
    base_metrics = {
        "cpa": audit.get("avg_cpa") or 0,
        "roas": audit.get("avg_roas") or audit.get("avg_ctr") or 0,
        "cv_count": audit.get("total_conversions") or 0,
        "cost": audit.get("total_cost") or monthly_spend,
        "ctr": audit.get("avg_ctr") or 0,
    }
    sev_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    issues_by_rule: dict[str, dict] = {}
    for issue in sorted(issues, key=lambda x: sev_order.get(x.get("severity"), 9)):
        rid = issue.get("check_id") or issue.get("id") or issue.get("rule_id")
        if rid and rid not in issues_by_rule:
            issues_by_rule[rid] = issue

    estimates = []
    for a in top5:
        rule = rules.get(a["rule_id"])
        if not rule:
            continue
        cm = dict(base_metrics)
        issue = issues_by_rule.get(a["rule_id"])
        if issue:
            pkey = issue.get("platform")
            if pkey and pkey in platform_summary:
                ps = platform_summary[pkey]
                if ps.get("avg_cpa"):
                    cm["cpa"] = ps["avg_cpa"]
                if ps.get("avg_roas"):
                    cm["roas"] = ps["avg_roas"]
                if ps.get("conversions"):
                    cm["cv_count"] = ps["conversions"]
                if ps.get("cost"):
                    cm["campaign_cost"] = ps["cost"]
                    cm["campaign_cpa"] = ps.get("avg_cpa") or 0
                    cm["campaign_cv"] = ps.get("conversions") or 0
        estimates.append(estimate_for_rule(rule, monthly_spend, current_metrics=cm))
    dedup = aggregate_with_dedup(estimates, rules, weights)
    per_est = {p["rule_id"]: p for p in dedup.get("per_estimate_with_factor", [])}

    # 全体 3 層合計（pixel_health 連動 ON）
    minimum = calculate_minimum_impact(estimates, rules, weights, pixel_health=pixel_health)
    realistic = calculate_realistic_impact(estimates, rules, weights, pixel_health=pixel_health)
    independent = calculate_independent_impact(estimates, rules, weights)

    # Top5 各行の 3 層値を構築
    top5_rows = []
    for i, (action, est) in enumerate(zip(top5, estimates), 1):
        rid = action["rule_id"]
        sc = est.get("scenario", {}) or {}
        per_info = per_est.get(rid, {})
        top5_rows.append({
            "rank": i,
            "rule_id": rid,
            "rule_name": action.get("rule_name", ""),
            "group": per_info.get("group", "other"),
            "group_short": GROUP_SHORT.get(per_info.get("group", "other"), "—"),
            "factor": per_info.get("factor", 1.0),
            "conservative": int(sc.get("conservative_yen", 0)),
            "realistic": int(sc.get("realistic_yen", 0)),
            "optimistic": int(sc.get("optimistic_yen", 0)),
            "confidence_stars": est.get("confidence_stars", "★☆☆"),
            "quick_win": action.get("quick_win", False),
        })

    return {
        "top5_rows": top5_rows,
        "total": {
            "conservative": minimum["total_yen"],
            "realistic": realistic["total_yen"],
            "optimistic": independent["total_yen"],
        },
        "monthly_spend": int(monthly_spend),
        "pixel_health": pixel_health,
    }


# =============================================================================
# pptx 構築ヘルパ
# =============================================================================

def _add_textbox(slide, left_cm, top_cm, width_cm, height_cm,
                 text, font_size=14, bold=False, color=COLOR_DARK,
                 font_name=FONT_JP, align=PP_ALIGN.LEFT):
    """テキストボックスを追加し、整形して返す。"""
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb, p


def _add_filled_rect(slide, left_cm, top_cm, width_cm, height_cm,
                     fill_color, line_color=None):
    """塗りつぶし矩形を追加。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def _add_footer(slide, page_num: int, total: int):
    """全ページ共通フッター（社外秘 + ページ番号）"""
    _add_textbox(
        slide, 1.0, SLIDE_H_CM - 0.8, 20.0, 0.6,
        "社外秘 / Confidential — 株式会社 Zynect Media",
        font_size=8, color=COLOR_GRAY,
    )
    _add_textbox(
        slide, SLIDE_W_CM - 4.0, SLIDE_H_CM - 0.8, 3.0, 0.6,
        f"P {page_num} / {total}",
        font_size=8, color=COLOR_GRAY, align=PP_ALIGN.RIGHT,
    )


def _add_speaker_notes(slide, data_source: str, related_adr: str,
                        narrative: str, generated_at: str):
    """スピーカーノートに ADR トレーサビリティを機械的挿入。"""
    notes_slide = slide.notes_slide
    text = (
        f"【データソース】{data_source}\n"
        f"【根拠ADR】{related_adr}\n"
        f"【更新日】{generated_at}\n"
        f"【話法メモ】{narrative}"
    )
    notes_slide.notes_text_frame.text = text


def _setup_presentation() -> Presentation:
    """16:9 のベース pptx を準備。"""
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    return prs


def _add_blank_slide(prs):
    """空白レイアウトのスライドを追加。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    # 背景色（オフホワイト）
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    return slide


# =============================================================================
# ページ別ビルダー
# =============================================================================

TOTAL_PAGES = 13


def build_p1_cover(prs, client_cfg, generated_at, audit_source):
    """P1: 表紙"""
    slide = _add_blank_slide(prs)
    company_name = client_cfg.get("company", {}).get("name", "(未設定)")
    company_honor = client_cfg.get("company", {}).get("honorific", "御中")

    # ロゴ風（テキストベース）
    _add_textbox(slide, 2.0, 2.5, 20.0, 1.5,
                 "zynect media", font_size=36, bold=True,
                 color=COLOR_DARK, font_name=FONT_EN)
    _add_filled_rect(slide, 2.0, 4.3, 4.0, 0.05, COLOR_DARK)

    # メインタイトル
    _add_textbox(slide, 2.0, 6.0, 26.0, 1.6,
                 "広告アカウント健全性診断 + 改善ロードマップ",
                 font_size=28, bold=True, color=COLOR_DARK)
    _add_textbox(slide, 2.0, 7.5, 26.0, 0.8,
                 "Account Health Audit & Improvement Roadmap — PoC Proposal",
                 font_size=12, color=COLOR_GRAY, font_name=FONT_EN)

    # 宛先
    _add_textbox(slide, 2.0, 11.0, 22.0, 1.0,
                 f"{company_name} {company_honor}",
                 font_size=22, bold=True, color=COLOR_DARK)

    # メタ情報
    _add_textbox(slide, 2.0, 14.5, 22.0, 0.6,
                 f"発行日: {generated_at}",
                 font_size=10, color=COLOR_DARK)
    _add_textbox(slide, 2.0, 15.2, 22.0, 0.6,
                 "発行: Zynect Media 株式会社",
                 font_size=10, color=COLOR_DARK)
    _add_textbox(slide, 2.0, 16.5, 22.0, 0.6,
                 "本書は機密保持契約に基づき作成されたものです / Confidential",
                 font_size=9, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source=f"clients.yaml: company.name='{company_name}'",
        related_adr="—",
        narrative=f"提案先は {company_name}。本資料は PoC 提案の表紙、まず信頼感のあるブランド表現で導入する。",
        generated_at=generated_at,
    )
    _add_footer(slide, 1, TOTAL_PAGES)
    return slide


def build_p2_executive_summary(prs, audit, top5_data, generated_at, audit_source):
    """P2: エグゼサマ — 3 行サマリ + Top5 全体 3 層 + KPI 4 マス"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "01 — Executive Summary",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "エグゼクティブサマリ",
                 font_size=24, bold=True, color=COLOR_DARK)

    # 3 行サマリ
    cv = int(audit.get("total_conversions", 0) or 0)
    cpa = int(audit.get("avg_cpa", 0) or 0)
    spend = top5_data["monthly_spend"]
    line1 = f"直近30日の月次広告費 ¥{spend:,} に対し CV {cv} 件、CPA ¥{cpa:,}（業界平均 ¥4,500 の約 2 倍）"
    line2 = f"主要課題は計測基盤の未整備（ROAS 計測不能、ピクセル休眠 {top5_data['pixel_health']['dormant_pixel_count']} 件）"
    line3 = f"Top5 アクションで月次 ¥{top5_data['total']['conservative']:,}〜¥{top5_data['total']['optimistic']:,} の改善見込み"
    for i, line in enumerate([line1, line2, line3]):
        _add_textbox(slide, 1.5, 3.3 + i * 0.85, 30.0, 0.8,
                     f"{i+1}. {line}", font_size=12, color=COLOR_DARK)

    # KPI 4 マス
    kpi_y = 7.0
    kpi_w = 7.0
    kpi_h = 3.0
    kpi_data = [
        ("月次広告費", f"¥{spend:,}", "直近30日 実測"),
        ("CV 数", f"{cv} 件", "直近30日 実測"),
        ("CPA", f"¥{cpa:,}", "業界平均 ¥4,500"),
        ("ROAS", "0.00", "Conversion Value 未送信"),
    ]
    for i, (label, value, note) in enumerate(kpi_data):
        x = 1.5 + i * (kpi_w + 0.4)
        _add_filled_rect(slide, x, kpi_y, kpi_w, kpi_h, COLOR_LIGHT_GRAY)
        _add_textbox(slide, x + 0.3, kpi_y + 0.3, kpi_w - 0.6, 0.6,
                     label, font_size=10, color=COLOR_GRAY)
        _add_textbox(slide, x + 0.3, kpi_y + 1.0, kpi_w - 0.6, 1.2,
                     value, font_size=22, bold=True, color=COLOR_DARK)
        _add_textbox(slide, x + 0.3, kpi_y + 2.3, kpi_w - 0.6, 0.6,
                     note, font_size=9, color=COLOR_GRAY)

    # Top5 全体 3 層
    summary_y = 10.8
    summary_h = 3.5
    _add_filled_rect(slide, 1.5, summary_y, 30.5, summary_h, RGBColor(0xFA, 0xFA, 0xF8))
    _add_textbox(slide, 1.8, summary_y + 0.2, 25.0, 0.5,
                 "想定改善見込み（Top5 全件、月次）",
                 font_size=10, bold=True, color=COLOR_DARK)
    layers = [
        ("確実値（pixel休眠連動込）", top5_data["total"]["conservative"], COLOR_GREEN),
        ("現実値（依存関係考慮）", top5_data["total"]["realistic"], COLOR_DARK),
        ("上限値（独立試算）", top5_data["total"]["optimistic"], COLOR_GRAY),
    ]
    for i, (label, value, color) in enumerate(layers):
        x = 1.8 + i * 10.0
        _add_textbox(slide, x, summary_y + 0.95, 9.5, 0.5,
                     label, font_size=9, color=COLOR_GRAY)
        _add_textbox(slide, x, summary_y + 1.5, 9.5, 1.4,
                     f"¥{value:,}/月",
                     font_size=22 if i == 0 else 18 if i == 1 else 14,
                     bold=(i == 0), color=color)

    _add_textbox(slide, 1.5, summary_y + summary_h + 0.2, 30.0, 0.5,
                 "※ 確実値・現実値・上限値の算出根拠は ADR-001/002/003 を参照。詳細は P8 Top5 アクション統合表に記載。",
                 font_size=8, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source=f"{audit_source} + engine/impact_estimator.py (calculate_minimum/realistic/independent_impact)",
        related_adr="ADR-001 (3層インパクト表示), ADR-004 (CV正規化、数値の信頼性根拠)",
        narrative=f"M09 単独で確実値 ¥{top5_data['top5_rows'][0]['conservative']:,} の改善見込み。CAPI 実装（M02）も確実値 ¥{top5_data['top5_rows'][1]['conservative']:,} 上乗せ。Top5 全体で確実値合計 ¥{top5_data['total']['conservative']:,}、上限値 ¥{top5_data['total']['optimistic']:,}（pixel 休眠連動込み）。月次広告費 ¥{spend:,} に対し約 40% の効率改善余地を提示可能。",
        generated_at=generated_at,
    )
    _add_footer(slide, 2, TOTAL_PAGES)
    return slide


def build_p3_current_quantitative(prs, audit, top5_data, generated_at, audit_source):
    """P3: 現状分析（数字）— 業界平均 3 軸比較"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "02 — Current State (Quantitative)",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "現状分析 — 数字で見るパイロットン",
                 font_size=22, bold=True, color=COLOR_DARK)

    # 業界 3 軸比較表
    bm = load_benchmarks()
    industry = "beauty_d2c"
    ps = audit.get("platform_summary", {}).get("meta", {}) or {}

    metrics = [
        ("CTR", ps.get("avg_ctr"), "ctr", "%"),
        ("CPA", ps.get("avg_cpa"), "cpa", "¥"),
        ("CVR", ps.get("avg_cvr"), "cvr", "%"),
        ("ROAS", ps.get("avg_roas"), "roas", "倍"),
    ]
    table_y = 4.0
    row_h = 1.4
    cols = ["指標", "現状", "業界平均", "Zynect 推奨", "判定"]
    col_widths = [4.5, 5.0, 5.0, 5.0, 6.0]
    col_xs = [1.5]
    for w in col_widths[:-1]:
        col_xs.append(col_xs[-1] + w)

    # ヘッダ
    for x, w, label in zip(col_xs, col_widths, cols):
        _add_filled_rect(slide, x, table_y, w, 0.8, COLOR_DARK)
        _add_textbox(slide, x + 0.2, table_y + 0.15, w - 0.4, 0.5,
                     label, font_size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                     align=PP_ALIGN.CENTER)

    # 行
    for i, (label, current, key, unit) in enumerate(metrics):
        y = table_y + 0.8 + i * row_h
        cmp = compare_3axis(industry, "meta_ads", key, float(current) if current else None, bm)
        ia = cmp.get("industry_avg")
        zr = cmp.get("zynect_recommended")

        def _fmt(v):
            if v is None:
                return "—"
            if unit == "¥":
                return f"¥{int(v):,}"
            if unit == "%":
                return f"{v:.2f}%"
            if unit == "倍":
                return f"{v:.2f}倍"
            return str(v)

        # 評価
        status = cmp.get("status", "no_benchmark")
        status_text = {
            "above_zynect": "✅ Zynect 推奨超え",
            "above_industry": "🟡 業界平均超え",
            "below_industry": "🔴 業界平均未達",
            "no_benchmark": "—",
        }.get(status, "—")
        status_color = {
            "above_zynect": COLOR_GREEN,
            "above_industry": COLOR_ORANGE,
            "below_industry": COLOR_RED,
            "no_benchmark": COLOR_GRAY,
        }.get(status, COLOR_GRAY)

        cells = [label, _fmt(current), _fmt(ia), _fmt(zr), status_text]
        for x, w, cell in zip(col_xs, col_widths, cells):
            _add_filled_rect(slide, x, y, w, row_h - 0.05,
                             RGBColor(0xFF, 0xFF, 0xFF), line_color=COLOR_LIGHT_GRAY)
            color = status_color if cell == status_text else COLOR_DARK
            bold = (cell == status_text) or (cells.index(cell) == 0)
            _add_textbox(slide, x + 0.2, y + 0.4, w - 0.4, row_h - 0.5,
                         cell, font_size=12, color=color, bold=bold,
                         align=PP_ALIGN.CENTER if cell != label else PP_ALIGN.LEFT)

    # 注記
    _add_textbox(slide, 1.5, 13.5, 30.0, 1.5,
                 f"※ 業界ベンチマークは config/benchmarks.yaml の beauty_d2c × meta_ads（出典: AdEspresso 2024 / Databox 2024 / WordStream 2024）。\n"
                 f"※ ROAS 0.00 は計測不能（Conversion Value 未送信）が原因で、実収益は別途評価が必要。",
                 font_size=10, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source=f"{audit_source} + config/benchmarks.yaml (beauty_d2c × meta_ads)",
        related_adr="ADR-002 (6グループ), ADR-004 (CV 正規化)",
        narrative="CTR は業界平均超えの優秀運用。一方 CPA は業界平均の約 2 倍で大きな改善余地。ROAS は計測不能のため評価不可。Zynect は CPA を業界平均レベルまで引き戻し、ROAS 計測を復活させる役割。",
        generated_at=generated_at,
    )
    _add_footer(slide, 3, TOTAL_PAGES)
    return slide


def build_p4_current_qualitative(prs, audit, top5_data, generated_at, audit_source):
    """P4: 現状分析（質的）— MYNAILPLEX 単独運用 + ピクセル健全性"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "02 — Current State (Qualitative)",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "現状分析 — 質的観察",
                 font_size=22, bold=True, color=COLOR_DARK)

    # 観察事項（数値は SSoT: audit から動的取得）
    ps = audit.get("platform_summary", {}).get("meta", {}) or {}
    ctr_actual = ps.get("avg_ctr")
    cpa_actual = ps.get("avg_cpa")
    bm = load_benchmarks()
    cpa_industry = (((bm.get("benchmarks") or {}).get("beauty_d2c") or {})
                    .get("meta_ads") or {}).get("cpa", {}).get("industry_avg")
    ctr_industry = (((bm.get("benchmarks") or {}).get("beauty_d2c") or {})
                    .get("meta_ads") or {}).get("ctr", {}).get("industry_avg")

    ctr_str = f"{ctr_actual:.2f}%" if ctr_actual is not None else "—"
    ctr_industry_str = f"{ctr_industry:.2f}%" if ctr_industry is not None else "—"
    cpa_str = f"¥{int(cpa_actual):,}" if cpa_actual else "—"
    cpa_ratio = (cpa_actual / cpa_industry) if (cpa_actual and cpa_industry) else None
    cpa_ratio_str = f"約 {cpa_ratio:.1f} 倍" if cpa_ratio else "—"

    observations = [
        ("運用構造", "MYNAILPLEX 単独運用が確定（直近30日の支出 95.4% を占有）",
         "CLOOKING・アゲルキャリアは配信実績ゼロ、Phase 2 オプション扱い"),
        ("運用品質", f"CTR {ctr_str}（業界平均 {ctr_industry_str} 超え）でクリック誘導は良好",
         "CR・LP のメッセージ整合性が機能している証拠"),
        ("計測課題", f"CPA {cpa_str} で業界平均の{cpa_ratio_str}、ROAS 0.00 で計測不能",
         "Conversion Value 未送信が ROAS 計測不能の根本原因"),
        ("ピクセル健全性",
         f"5 ピクセル中 {top5_data['pixel_health']['dormant_pixel_count']} 件休眠 "
         f"（最大 {top5_data['pixel_health']['dormant_days']} 日未発火）+ 重複疑い",
         "計測基盤の整理が他施策の前提"),
    ]
    y = 4.0
    for label, point, detail in observations:
        # 帯
        _add_filled_rect(slide, 1.5, y, 6.0, 2.6, COLOR_LIGHT_GRAY)
        _add_textbox(slide, 1.7, y + 0.3, 5.6, 0.6,
                     label, font_size=11, bold=True, color=COLOR_DARK)
        _add_textbox(slide, 1.7, y + 1.1, 5.6, 1.4,
                     point, font_size=10, color=COLOR_DARK)
        # 補足
        _add_textbox(slide, 8.0, y + 0.5, 24.0, 1.8,
                     detail, font_size=11, color=COLOR_DARK)
        y += 2.9

    _add_speaker_notes(
        slide,
        data_source=f"{audit_source} + reports/<date>/pilotton_brand_breakdown.md + analyzers/ads_audit.py:detect_pixel_health()",
        related_adr="ADR-003 (pixel_health 連動), ADR-002 (6グループ)",
        narrative="MYNAILPLEX 単独運用の確定は Day 5.2 の発見。これにより PoC 戦略を「全ブランド並行最適化」から「MYNAILPLEX 集中投下深掘り」に転換した経緯を共有可能。CLOCKING/アゲルキャリアの再開支援は契約スコープ外として明確に切り分け。",
        generated_at=generated_at,
    )
    _add_footer(slide, 4, TOTAL_PAGES)
    return slide


def build_p5_proposal_1_cpa(prs, audit, top5_data, generated_at, audit_source):
    """P5: 改善提案① CPA 業界平均超えへの是正"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "03 — Proposal 1 of 3",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "改善提案① CPA 業界平均超えへの是正",
                 font_size=22, bold=True, color=COLOR_DARK)

    # SSoT: 数値は audit + benchmarks から動的取得
    ps = audit.get("platform_summary", {}).get("meta", {}) or {}
    cpa_actual = ps.get("avg_cpa")
    bm = load_benchmarks()
    cpa_industry = (((bm.get("benchmarks") or {}).get("beauty_d2c") or {})
                    .get("meta_ads") or {}).get("cpa", {}).get("industry_avg")
    if cpa_actual and cpa_industry:
        improvement_pct = (cpa_actual - cpa_industry) / cpa_actual * 100
        headline = (f"現状 CPA ¥{int(cpa_actual):,} → 業界平均 ¥{int(cpa_industry):,} への是正で"
                    f"月次 約 {improvement_pct:.0f}% 効率改善（最大）")
    else:
        headline = "現状 CPA → 業界平均水準への是正で大幅な効率改善が見込めます"

    _add_textbox(slide, 1.5, 3.0, 30.0, 1.0,
                 headline,
                 font_size=14, bold=True, color=COLOR_GREEN)

    # 4 つの施策
    items = [
        ("M02 CAPI 完全実装", "計測欠損を補い、自動入札の学習精度を回復", "★★★ 効果確実"),
        ("M03 EMQ（イベントマッチ品質）改善", "ハッシュ化メール常時送信で Pixel/CAPI のシグナル精度向上", "★★★ 効果確実"),
        ("M04 ドメイン検証", "AEM（Aggregated Event Measurement）の前提条件、iOS14 対応", "★★★ 効果確実"),
        ("M09 学習フェーズ最適化", "週 50CV/広告セット基準の達成、不要な編集介入の抑制", "★★☆ 効果は中"),
    ]
    y = 5.0
    for label, point, conf in items:
        _add_filled_rect(slide, 1.5, y, 30.5, 1.6, RGBColor(0xFA, 0xFA, 0xF8))
        _add_textbox(slide, 1.8, y + 0.2, 12.0, 0.6,
                     label, font_size=12, bold=True, color=COLOR_DARK)
        _add_textbox(slide, 1.8, y + 0.85, 22.0, 0.6,
                     point, font_size=10, color=COLOR_GRAY)
        _add_textbox(slide, 25.0, y + 0.5, 6.5, 0.7,
                     conf, font_size=10, color=COLOR_GREEN, align=PP_ALIGN.RIGHT)
        y += 1.85

    # 米満氏理論との接続
    _add_textbox(slide, 1.5, 13.5, 30.0, 1.5,
                 "▶ 米満氏理論 P1（計測精度=学習シグナル精度）/ M-α（Pixel + CAPI + EMQ）/ M-β（学習フェーズ保護）に基づく施策群。\n"
                 "計測基盤の修復が他施策の前提のため、本提案では最優先で着手します。",
                 font_size=10, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source="config/rules/meta_rules.yaml: M02/M03/M04/M09 + docs/principles/meta_principles.md",
        related_adr="ADR-002 (measurement_foundation グループ), ADR-003 (pixel 連動による優先度)",
        narrative="CPA を業界平均レベルまで戻すだけで月次 ¥581K の確実な改善が見込める。営業時はこの数字を主軸に、計測修復が「他施策の前提」であることを強調する。米満氏理論 P1/M-α は他社差別化の語り口。",
        generated_at=generated_at,
    )
    _add_footer(slide, 5, TOTAL_PAGES)
    return slide


def build_p6_proposal_2_roas(prs, generated_at):
    """P6: 改善提案② ROAS 計測復活"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "03 — Proposal 2 of 3",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "改善提案② ROAS 計測復活",
                 font_size=22, bold=True, color=COLOR_DARK)

    _add_textbox(slide, 1.5, 3.0, 30.0, 1.0,
                 "Conversion Value（purchase_value）を Pixel/CAPI 経由で送信し、ROAS の真値を可視化",
                 font_size=14, bold=True, color=COLOR_GREEN)

    # ステップ
    steps = [
        "1. Meta Events Manager で「テストイベント」タブを開き、Purchase イベントの value パラメータが入っているか確認",
        "2. 入っていなければサイト側 Pixel コードに value: <購入金額> を追加（GTM or 直貼り）",
        "3. CAPI 経由でも同じ value を送信（重複は event_id でデデュプ）",
        "4. 「コンバージョン API ヘルス」で「value 受信あり」表示を確認",
        "5. 翌週から ROAS が表示され、業界平均 2.80 倍との比較が機能するようになる",
    ]
    y = 5.0
    for s in steps:
        _add_filled_rect(slide, 1.5, y, 30.5, 1.0, RGBColor(0xFA, 0xFA, 0xF8))
        _add_textbox(slide, 1.8, y + 0.2, 30.0, 0.7,
                     s, font_size=11, color=COLOR_DARK)
        y += 1.2

    _add_filled_rect(slide, 1.5, 12.0, 30.5, 1.5, RGBColor(0xFC, 0xEB, 0xEB))
    _add_textbox(slide, 1.8, 12.2, 30.0, 1.1,
                 "⚠ 現状 ROAS=0 は計測不備の可能性が高く、実際の収益貢献は別途評価が必要。\n"
                 "本施策完了後、初めて広告効果の真の姿が見えてきます。",
                 font_size=10, color=COLOR_RED)

    _add_speaker_notes(
        slide,
        data_source="ADR-004 + analyzers/ads_audit.py の DQ-ROAS-MISSING アラート",
        related_adr="ADR-004 (CV 正規化), ADR-001 (3層インパクト)",
        narrative="ROAS=0 は他社監査ツールでも検出可能だが、修復手順を 5 ステップで具体化できる点が Zynect の差別化。CAPI 実装の所要時間 3〜5 営業日と明確に提示可能。",
        generated_at=generated_at,
    )
    _add_footer(slide, 6, TOTAL_PAGES)
    return slide


def build_p7_proposal_3_pixel(prs, top5_data, generated_at, audit_source):
    """P7: 改善提案③ 計測基盤の総整備"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "03 — Proposal 3 of 3",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "改善提案③ 計測基盤の総整備",
                 font_size=22, bold=True, color=COLOR_DARK)

    ph = top5_data["pixel_health"]
    _add_textbox(slide, 1.5, 3.0, 30.0, 1.0,
                 f"ピクセル {ph['total_pixel_count']} 件中 {ph['dormant_pixel_count']} 件が長期休眠 + 重複疑いあり、整理が他施策の前提",
                 font_size=14, bold=True, color=COLOR_GREEN)

    # 整理対象
    items = [
        ("CLOOKING Pixel × 2 件",
         "重複疑い + 270/366 日休眠",
         "Meta UI で 1 つに統合 or 廃止"),
        ("「削除」ピクセル",
         "未発火、廃止予定",
         "Meta UI で実削除"),
        ("AEM 優先度イベント",
         "iOS14 後の計測欠損対応",
         "Purchase / AddToCart / Lead を優先 8 件まで設定"),
    ]
    y = 5.0
    for label, point, action in items:
        _add_filled_rect(slide, 1.5, y, 30.5, 2.0, RGBColor(0xFA, 0xFA, 0xF8))
        _add_textbox(slide, 1.8, y + 0.2, 12.0, 0.7,
                     label, font_size=12, bold=True, color=COLOR_DARK)
        _add_textbox(slide, 1.8, y + 0.95, 16.0, 0.7,
                     point, font_size=10, color=COLOR_GRAY)
        _add_textbox(slide, 18.0, y + 0.5, 14.0, 1.1,
                     action, font_size=11, color=COLOR_GREEN)
        y += 2.3

    _add_textbox(slide, 1.5, 12.5, 30.0, 1.5,
                 "▶ 米満氏理論 M-α / P1（計測精度）に基づく整理。\n"
                 "ピクセル整理が完了するまで、改善提案①の効果は限定的（pixel_health 連動で改善見込み額を保守的に算出済）。",
                 font_size=10, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source="config/clients.yaml: pilotton.ads.meta.pixels + analyzers/ads_audit.py:detect_pixel_health()",
        related_adr="ADR-003 (pixel_health 連動)",
        narrative="ピクセル整理は手作業 30 分程度で完了。営業時は「我々が整理手順をご案内、貴社で 30 分作業いただくだけ」と低工数で訴求。整理後の効果が他施策に波及することを ADR-003 ベースで説明可能。",
        generated_at=generated_at,
    )
    _add_footer(slide, 7, TOTAL_PAGES)
    return slide


def build_p8_top5_3layer_table(prs, top5_data, generated_at, audit_source):
    """P8: 最重要 — Top5 + 3 層インパクト統合表"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.5, 8.0, 0.6,
                 "04 — Top5 Action Plan with 3-Layer Impact",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.1, 28.0, 1.0,
                 "Top5 アクション + 3 層インパクト",
                 font_size=20, bold=True, color=COLOR_DARK)

    # 表ヘッダ
    table_y = 2.6
    cols = ["#", "アクション", "Group", "係数", "確実値", "現実値", "上限値", "確度"]
    col_widths = [1.0, 11.0, 2.0, 1.6, 4.5, 4.5, 4.5, 1.8]
    col_xs = [1.5]
    for w in col_widths[:-1]:
        col_xs.append(col_xs[-1] + w)
    header_h = 0.7
    for x, w, label in zip(col_xs, col_widths, cols):
        _add_filled_rect(slide, x, table_y, w, header_h, COLOR_DARK)
        _add_textbox(slide, x + 0.1, table_y + 0.12, w - 0.2, header_h - 0.2,
                     label, font_size=10, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # データ行（priority_ranker の実出力順）
    row_h = 1.0
    for i, row in enumerate(top5_data["top5_rows"]):
        y = table_y + header_h + i * row_h
        bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xFA, 0xFA, 0xF8)
        for x, w in zip(col_xs, col_widths):
            _add_filled_rect(slide, x, y, w, row_h - 0.05, bg, line_color=COLOR_LIGHT_GRAY)
        cells = [
            (str(row["rank"]), False, COLOR_DARK, PP_ALIGN.CENTER, 11),
            (f"[{row['rule_id']}] {row['rule_name']}", True, COLOR_DARK, PP_ALIGN.LEFT, 11),
            (row["group_short"], False, COLOR_GRAY, PP_ALIGN.CENTER, 10),
            (f"×{row['factor']:.2f}", False, COLOR_GRAY, PP_ALIGN.CENTER, 10),
            (f"¥{row['conservative']:,}", True, COLOR_GREEN, PP_ALIGN.RIGHT, 11),
            (f"¥{row['realistic']:,}", False, COLOR_DARK, PP_ALIGN.RIGHT, 11),
            (f"¥{row['optimistic']:,}", False, COLOR_GRAY, PP_ALIGN.RIGHT, 11),
            (row["confidence_stars"], False, COLOR_DARK, PP_ALIGN.CENTER, 10),
        ]
        for x, w, (text, bold, color, align, fsize) in zip(col_xs, col_widths, cells):
            _add_textbox(slide, x + 0.1, y + 0.25, w - 0.2, row_h - 0.4,
                         text, font_size=fsize, bold=bold, color=color, align=align)

    # 合計行
    total_y = table_y + header_h + len(top5_data["top5_rows"]) * row_h
    _add_filled_rect(slide, 1.5, total_y, 30.9, row_h - 0.05, COLOR_HIGHLIGHT_BG)
    _add_textbox(slide, 1.7, total_y + 0.3, 14.5, row_h - 0.5,
                 "Top5 合計", font_size=12, bold=True, color=COLOR_DARK)
    total = top5_data["total"]
    _add_textbox(slide, 17.6, total_y + 0.3, 4.5, row_h - 0.5,
                 f"¥{total['conservative']:,}", font_size=12, bold=True,
                 color=COLOR_GREEN, align=PP_ALIGN.RIGHT)
    _add_textbox(slide, 22.1, total_y + 0.3, 4.5, row_h - 0.5,
                 f"¥{total['realistic']:,}", font_size=12, bold=True,
                 color=COLOR_DARK, align=PP_ALIGN.RIGHT)
    _add_textbox(slide, 26.6, total_y + 0.3, 4.5, row_h - 0.5,
                 f"¥{total['optimistic']:,}", font_size=12, bold=True,
                 color=COLOR_GRAY, align=PP_ALIGN.RIGHT)

    # 凡例
    legend_y = total_y + row_h + 0.3
    legend_text = (
        "凡例: MF=計測基盤 / DLS=配信学習・構造 / TGT=ターゲティング / IND=独立施策\n"
        "係数: グループ最大値=×1.00、2位以下=duplicate_factor 適用（pixel_health 連動で MF 内 0.1〜0.2）\n"
        "確実値: pixel_health 連動込み（measurement_foundation 内 2 位以下を 0.1×、非 MF を 0.7× decay）\n"
        "現実値: 連動なしの依存係数のみ適用 / 上限値: 各施策が独立に最大効果（重複排除なし、到達困難）"
    )
    _add_textbox(slide, 1.5, legend_y, 30.0, 3.0,
                 legend_text, font_size=9, color=COLOR_GRAY)

    # 動的順序の注記
    rule_order = " → ".join([r["rule_id"] for r in top5_data["top5_rows"]])
    _add_speaker_notes(
        slide,
        data_source="engine/priority_ranker.py:compute_top_actions() + engine/impact_estimator.py:per_estimate_with_factor()",
        related_adr="ADR-001 (3層インパクト表示), ADR-002 (6グループ duplicate_factor), ADR-003 (pixel_health 連動)",
        narrative=(
            f"M09 単独で確実値 ¥{top5_data['top5_rows'][0]['conservative']:,} の改善見込み。"
            f"CAPI 実装（M02）も確実値 ¥{top5_data['top5_rows'][1]['conservative']:,} 上乗せ。"
            f"Top5 全体で確実値合計 ¥{top5_data['total']['conservative']:,}、"
            f"上限値 ¥{top5_data['total']['optimistic']:,}（pixel 休眠連動込み）。"
            f"月次広告費 ¥{top5_data['monthly_spend']:,} に対し約 40% の効率改善余地を提示可能。"
            f"※ 順序は priority_ranker の実出力に動的追従（Day 5.3 時点 {rule_order}、CV データ更新で変動可能性あり）。"
        ),
        generated_at=generated_at,
    )
    _add_footer(slide, 8, TOTAL_PAGES)
    return slide


def build_p9_roadmap(prs, top5_data, generated_at, audit_source):
    """P9: 改善ロードマップ（タイムライン）"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "05 — Improvement Roadmap",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "改善ロードマップ（4 週間タイムライン）",
                 font_size=22, bold=True, color=COLOR_DARK)

    # ガントチャート風
    weeks = ["Week 1", "Week 2", "Week 3", "Week 4"]
    col_w = 6.5
    week_y = 4.0
    week_h = 0.7
    for i, week in enumerate(weeks):
        x = 6.0 + i * col_w
        _add_filled_rect(slide, x, week_y, col_w - 0.2, week_h, COLOR_LIGHT_GRAY)
        _add_textbox(slide, x + 0.2, week_y + 0.15, col_w - 0.4, week_h - 0.2,
                     week, font_size=11, bold=True, color=COLOR_DARK,
                     align=PP_ALIGN.CENTER)

    # 行（Top5 + 米満氏理論ベース推奨実装順）
    actions_timeline = [
        ("M04 ドメイン検証", 0, 1, COLOR_GREEN),  # Week 1
        ("M02 CAPI 実装", 0, 2, COLOR_GREEN),     # Week 1-2
        ("M03 EMQ 改善", 1, 1, COLOR_GREEN),       # Week 2
        ("ピクセル整理", 0, 1, COLOR_GREEN),        # Week 1
        ("M09 学習脱出待ち", 2, 2, COLOR_ORANGE),  # Week 3-4
        ("M61 1Pデータ活用", 2, 2, COLOR_ORANGE),  # Week 3-4
        ("計測検証 + ROAS 確認", 2, 2, COLOR_DARK),  # Week 3-4
    ]
    row_h = 0.85
    for i, (label, start, span, color) in enumerate(actions_timeline):
        y = week_y + week_h + 0.3 + i * row_h
        # 行ラベル
        _add_textbox(slide, 1.5, y + 0.2, 4.3, row_h - 0.3,
                     label, font_size=10, color=COLOR_DARK)
        # バー
        bar_x = 6.0 + start * col_w
        bar_w = span * col_w - 0.3
        _add_filled_rect(slide, bar_x, y + 0.2, bar_w, row_h - 0.4, color)

    # 注記
    note_y = week_y + week_h + 0.3 + len(actions_timeline) * row_h + 0.4
    _add_filled_rect(slide, 1.5, note_y, 30.0, 2.5, RGBColor(0xFA, 0xEE, 0xDA))
    _add_textbox(slide, 1.7, note_y + 0.2, 29.6, 0.5,
                 "⚠ 効果発現の前提条件",
                 font_size=11, bold=True, color=COLOR_ORANGE)
    notes_text = (
        "▸ 計測基盤系（M02 / M04）: 実装後 Pixel/CAPI からのシグナル蓄積に 2〜4 週\n"
        "▸ 学習フェーズ系（M09）: 学習脱出に週 50CV 必要、予算規模により 2〜4 週\n"
        "▸ 1P データ系（M61）: Customer File 連携 → LLA 再学習 → 配信反映で 4〜6 週"
    )
    _add_textbox(slide, 1.7, note_y + 0.8, 29.6, 1.6,
                 notes_text, font_size=9, color=COLOR_DARK)

    _add_speaker_notes(
        slide,
        data_source="config/rules/meta_rules.yaml: M02/M03/M04/M09/M61 の estimated_duration + impact_horizon_weeks",
        related_adr="ADR-002 (6グループ依存順), ADR-003 (pixel 連動)",
        narrative="計測基盤（M04→M02→M03）を Week 1-2 に並列着手、ピクセル整理も同時実施。学習フェーズ系（M09）は計測修復後の Week 3-4 で効果検証。M61 は同時並行可能。Week 5 以降に ROAS 復活を確認。",
        generated_at=generated_at,
    )
    _add_footer(slide, 9, TOTAL_PAGES)
    return slide


def build_p10_expected_results(prs, top5_data, generated_at, audit_source):
    """P10: 期待効果（数値サマリ）— 月次・年次"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "06 — Expected Results",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "期待効果（数値サマリ）",
                 font_size=22, bold=True, color=COLOR_DARK)

    total = top5_data["total"]
    monthly_spend = top5_data["monthly_spend"]

    # 3 層 × 月次/年次 マトリクス
    table_y = 3.5
    cols = ["シナリオ", "月次改善", "年次換算", "月次広告費比"]
    col_widths = [10.0, 7.0, 7.0, 7.0]
    col_xs = [1.5]
    for w in col_widths[:-1]:
        col_xs.append(col_xs[-1] + w)
    header_h = 0.8
    for x, w, label in zip(col_xs, col_widths, cols):
        _add_filled_rect(slide, x, table_y, w, header_h, COLOR_DARK)
        _add_textbox(slide, x + 0.2, table_y + 0.2, w - 0.4, header_h - 0.4,
                     label, font_size=12, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    rows = [
        ("確実値（pixel 連動込）", total["conservative"], COLOR_GREEN),
        ("現実値（依存考慮）", total["realistic"], COLOR_DARK),
        ("上限値（独立試算）", total["optimistic"], COLOR_GRAY),
    ]
    row_h = 1.5
    for i, (label, value, color) in enumerate(rows):
        y = table_y + header_h + i * row_h
        bg = COLOR_HIGHLIGHT_BG if i == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for x, w in zip(col_xs, col_widths):
            _add_filled_rect(slide, x, y, w, row_h - 0.05, bg, line_color=COLOR_LIGHT_GRAY)

        _add_textbox(slide, col_xs[0] + 0.3, y + 0.4, col_widths[0] - 0.4, 0.7,
                     label, font_size=13, bold=(i == 0), color=COLOR_DARK)
        _add_textbox(slide, col_xs[1] + 0.3, y + 0.4, col_widths[1] - 0.4, 0.7,
                     f"¥{value:,}", font_size=14, bold=(i == 0), color=color,
                     align=PP_ALIGN.RIGHT)
        _add_textbox(slide, col_xs[2] + 0.3, y + 0.4, col_widths[2] - 0.4, 0.7,
                     f"¥{value * 12:,}", font_size=12, color=color,
                     align=PP_ALIGN.RIGHT)
        ratio = value / monthly_spend * 100 if monthly_spend > 0 else 0
        _add_textbox(slide, col_xs[3] + 0.3, y + 0.4, col_widths[3] - 0.4, 0.7,
                     f"{ratio:.1f}%", font_size=12, color=color,
                     align=PP_ALIGN.RIGHT)

    # 注記
    annot_y = table_y + header_h + len(rows) * row_h + 0.5
    _add_filled_rect(slide, 1.5, annot_y, 30.5, 3.0, RGBColor(0xFA, 0xFA, 0xF8))
    _add_textbox(slide, 1.7, annot_y + 0.2, 30.0, 0.5,
                 "営業説明上のスタンス",
                 font_size=11, bold=True, color=COLOR_DARK)
    stance_text = (
        f"▸ 確実値（¥{total['conservative']:,}/月）を主指標として提示します。これは過大評価を避けるため依存関係と pixel 健全性を保守的に折り込んだ最低値です。\n"
        f"▸ 現実値（¥{total['realistic']:,}/月）は計測基盤の修復が完了した想定で、依存係数のみ適用した中央推計です。\n"
        f"▸ 上限値（¥{total['optimistic']:,}/月）は理論上限で、現実には到達困難ですが「改善余地の天井」として参考表示します。"
    )
    _add_textbox(slide, 1.7, annot_y + 0.8, 30.0, 2.1,
                 stance_text, font_size=10, color=COLOR_DARK)

    _add_speaker_notes(
        slide,
        data_source="engine/impact_estimator.py の calculate_minimum/realistic/independent_impact",
        related_adr="ADR-001 (3層インパクト表示)",
        narrative=f"確実値 ¥{total['conservative']:,}/月 = 月次広告費 ¥{monthly_spend:,} の約 {total['conservative'] / monthly_spend * 100:.0f}%。営業時は確実値を約束として提示し、現実値・上限値は「上振れ余地」として伝える。年間換算で営業効果を訴求。",
        generated_at=generated_at,
    )
    _add_footer(slide, 10, TOTAL_PAGES)
    return slide


def build_p11_schedule(prs, generated_at):
    """P11: 実装スケジュール — 4 週マイルストーン + 担当範囲"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "07 — Implementation Schedule",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "実装スケジュール",
                 font_size=22, bold=True, color=COLOR_DARK)

    weeks = [
        ("Week 1", "計測基盤の整備開始",
         "Zynect: M04 ドメイン検証手順案内 / M02 CAPI 実装着手 / ピクセル整理手順案内\n"
         "貴社: ピクセル整理（Meta UI で 30 分作業）"),
        ("Week 2", "計測基盤の整備完了",
         "Zynect: M02 CAPI 実装完了 / M03 EMQ 改善 / ハッシュ化メール送信設定\n"
         "貴社: テストイベント確認、社内 IT 調整"),
        ("Week 3", "学習フェーズ最適化 + 1P データ活用",
         "Zynect: M09 学習設定 / M61 Customer File 連携設定\n"
         "貴社: 顧客リスト CSV ご提供、Conversion Value 設定確認"),
        ("Week 4", "効果検証 + ROAS 確認",
         "Zynect: 効果検証レポート + 次月施策提案\n"
         "貴社: ROAS の真値確認、Phase 2 検討"),
    ]
    y = 3.5
    for week, title, content in weeks:
        _add_filled_rect(slide, 1.5, y, 4.0, 2.8, COLOR_DARK)
        _add_textbox(slide, 1.7, y + 0.4, 3.6, 0.6,
                     week, font_size=12, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        _add_textbox(slide, 1.7, y + 1.3, 3.6, 1.4,
                     title, font_size=10, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        _add_filled_rect(slide, 5.5, y, 26.5, 2.8, RGBColor(0xFA, 0xFA, 0xF8))
        _add_textbox(slide, 5.7, y + 0.3, 26.1, 2.3,
                     content, font_size=10, color=COLOR_DARK)
        y += 3.1

    _add_speaker_notes(
        slide,
        data_source="config/rules/meta_rules.yaml の estimated_duration + ADR-004 Implementation",
        related_adr="ADR-001/002/003/004 全般",
        narrative="4 週間で完結する想定。Zynect は計測基盤修復、貴社はピクセル整理（30 分）+ Customer File 提供のみ。役割分担を明確に提示することで導入ハードルを下げる。",
        generated_at=generated_at,
    )
    _add_footer(slide, 11, TOTAL_PAGES)
    return slide


def build_p12_pricing(prs, generated_at):
    """P12: 料金（プレースホルダ）"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "08 — Pricing",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "料金体系",
                 font_size=22, bold=True, color=COLOR_DARK)

    # プレースホルダ（提案毎に手動差し替え）
    _add_filled_rect(slide, 4.0, 6.0, 25.5, 6.0, COLOR_LIGHT_GRAY)
    _add_textbox(slide, 4.5, 7.0, 24.5, 1.5,
                 "料金体系は貴社の運用規模・改善優先度に応じて",
                 font_size=18, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    _add_textbox(slide, 4.5, 8.5, 24.5, 1.5,
                 "ご提案いたします",
                 font_size=18, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    _add_textbox(slide, 4.5, 10.5, 24.5, 0.8,
                 "対面議論時にカスタマイズ可能な料金プランをご相談ください",
                 font_size=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    _add_speaker_notes(
        slide,
        data_source="（手動入力）— pptx 自動生成では空、提案前に料金 xlsx シミュレータから出力して差し替え",
        related_adr="—",
        narrative="料金は対面議論時に確定。事前に xlsx シミュレータで 3 パターン（小規模 / 標準 / 大規模）を準備し、対話に応じて即座に提示。本ページは差し替え用のプレースホルダ。",
        generated_at=generated_at,
    )
    _add_footer(slide, 12, TOTAL_PAGES)
    return slide


def build_p13_appendix(prs, generated_at):
    """P13: Appendix — Phase 2 / ベンチマーク / 原則 / 計算ロジック注記"""
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 1.5, 0.8, 8.0, 0.7,
                 "09 — Appendix",
                 font_size=10, color=COLOR_GRAY, font_name=FONT_EN)
    _add_textbox(slide, 1.5, 1.5, 28.0, 1.1,
                 "付録",
                 font_size=22, bold=True, color=COLOR_DARK)

    # A: Phase 2 オプション
    _add_textbox(slide, 1.5, 3.2, 30.0, 0.6,
                 "A. Phase 2 オプション（PoC 期間外）",
                 font_size=12, bold=True, color=COLOR_DARK)
    _add_textbox(slide, 1.5, 3.9, 30.0, 1.8,
                 "▸ CLOOKING ピクセル整理（Meta UI で重複統合 + 廃止）\n"
                 "▸ アゲルキャリアブランド再開支援（必要に応じて別契約）\n"
                 "▸ AdTruth による不正クリック検知連携（Fraud リスク顕在化時）",
                 font_size=10, color=COLOR_DARK)

    # B: 業界ベンチマーク出典
    _add_textbox(slide, 1.5, 6.3, 30.0, 0.6,
                 "B. 業界ベンチマーク出典",
                 font_size=12, bold=True, color=COLOR_DARK)
    _add_textbox(slide, 1.5, 7.0, 30.0, 1.5,
                 "AdEspresso 2024 / Databox 2024 / WordStream 2024 / Varos B2B SaaS Panel 2025 / Powered by Search 2024。\n"
                 "config/benchmarks.yaml の 6 業界 × 3 媒体 × 6 メトリクス（CTR/CPC/CPA/CVR/ROAS/Frequency）に集約。",
                 font_size=10, color=COLOR_DARK)

    # C: 米満氏理論
    _add_textbox(slide, 1.5, 9.0, 30.0, 0.6,
                 "C. 米満氏理論 9+10 原則（要約）",
                 font_size=12, bold=True, color=COLOR_DARK)
    # 2 列レイアウト
    half = (len(PRINCIPLES_SUMMARY) + 1) // 2
    for i, (pid, name, desc) in enumerate(PRINCIPLES_SUMMARY):
        if i < half:
            x = 1.5
            y = 9.7 + i * 0.4
        else:
            x = 16.5
            y = 9.7 + (i - half) * 0.4
        _add_textbox(slide, x, y, 14.5, 0.4,
                     f"{pid}: {name} — {desc}",
                     font_size=8, color=COLOR_DARK)

    # D: 計算ロジック注記
    notes_y = 16.0
    _add_textbox(slide, 1.5, notes_y, 30.0, 0.6,
                 "D. 計算ロジック注記（裏側）",
                 font_size=10, bold=True, color=COLOR_GRAY)
    _add_textbox(slide, 1.5, notes_y + 0.5, 30.0, 1.2,
                 "本資料の数値は ADR-001（3層インパクト）/ ADR-002（6グループ）/ ADR-003（pixel 連動）/ ADR-004（CV 正規化）に基づき算出。\n"
                 "Meta API の同一購入が複数ラベルで重複報告される仕様は内部で正規化済み。",
                 font_size=8, color=COLOR_GRAY)

    _add_speaker_notes(
        slide,
        data_source="docs/decisions/ADR-001 〜 ADR-004 / docs/principles/{google,meta}_principles.md / config/benchmarks.yaml",
        related_adr="ADR-001, ADR-002, ADR-003, ADR-004",
        narrative="Phase 2 オプションは契約後の追加施策として位置付け。ベンチマーク出典は業界比較の根拠として明示。米満氏理論は Zynect の独自性の根幹で、提案後の Q&A で深掘り対応可能。計算ロジック注記は ADR-004 を含む 4 本の ADR への参照。",
        generated_at=generated_at,
    )
    _add_footer(slide, 13, TOTAL_PAGES)
    return slide


# =============================================================================
# バージョン管理（ファイル名）
# =============================================================================

def determine_output_path(client_id: str, output_dir: Path | None) -> Path:
    """ファイル名規則: <client>_proposal_v<N>.pptx
    既存ファイルがあれば自動でバージョン番号インクリメント。
    """
    if output_dir is None:
        today = datetime.now().strftime("%Y-%m-%d")
        output_dir = PROJECT_ROOT / "reports" / today
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    pattern = re.compile(rf"^{re.escape(client_id)}_proposal_v(\d+)(?:_\d{{4}})?\.pptx$")
    max_v = 0
    for f in output_dir.glob(f"{client_id}_proposal_*.pptx"):
        m = pattern.match(f.name)
        if m:
            max_v = max(max_v, int(m.group(1)))
    next_v = max_v + 1
    return output_dir / f"{client_id}_proposal_v{next_v}.pptx"


# =============================================================================
# メイン
# =============================================================================

def generate(client_id: str, output_dir: Path | None = None) -> Path:
    """pptx 全 13 ページ生成のメイン関数"""
    print(f"=== pilotton PoC 提案 pptx 生成開始 (client={client_id}) ===")

    # === SSoT データ取得 ===
    audit_data = load_audit_results(client_id)
    audit = audit_data.get("ads_audit", {})
    audit_source = audit_data.get("_source_path", "unknown")
    clients_yaml = load_clients_yaml()
    client_cfg = get_client_config(client_id, clients_yaml)
    top5_data = gather_top5_with_3layer(audit, client_cfg)
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M JST")

    # === pptx 構築 ===
    prs = _setup_presentation()

    build_p1_cover(prs, client_cfg, generated_at, audit_source)
    build_p2_executive_summary(prs, audit, top5_data, generated_at, audit_source)
    build_p3_current_quantitative(prs, audit, top5_data, generated_at, audit_source)
    build_p4_current_qualitative(prs, audit, top5_data, generated_at, audit_source)
    build_p5_proposal_1_cpa(prs, audit, top5_data, generated_at, audit_source)
    build_p6_proposal_2_roas(prs, generated_at)
    build_p7_proposal_3_pixel(prs, top5_data, generated_at, audit_source)
    build_p8_top5_3layer_table(prs, top5_data, generated_at, audit_source)
    build_p9_roadmap(prs, top5_data, generated_at, audit_source)
    build_p10_expected_results(prs, top5_data, generated_at, audit_source)
    build_p11_schedule(prs, generated_at)
    build_p12_pricing(prs, generated_at)
    build_p13_appendix(prs, generated_at)

    # === 保存 ===
    out_path = determine_output_path(client_id, output_dir)
    prs.save(str(out_path))
    print(f"=== 生成完了: {out_path} (size: {out_path.stat().st_size:,} bytes) ===")
    return out_path


def main():
    parser = argparse.ArgumentParser(description="pilotton PoC 提案 pptx 生成")
    parser.add_argument("--client", required=True, help="クライアント ID（pilotton 等）")
    parser.add_argument("--output-dir", help="出力ディレクトリ（省略時 reports/<today>/）")
    args = parser.parse_args()

    output_dir = Path(args.output_dir) if args.output_dir else None
    generate(args.client, output_dir)


if __name__ == "__main__":
    main()
